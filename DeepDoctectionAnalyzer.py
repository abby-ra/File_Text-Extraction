"""
Unified Document Extraction and Classification System
Combines OCR text extraction with AI-powered document type classification
Supports 15+ file formats with automatic document categorization
"""

import os
import sys
from pathlib import Path
import json
import cv2
import numpy as np
import fitz  # PyMuPDF
from docx import Document as docx_Document
from pptx import Presentation
import openpyxl
import torch
from PIL import Image
from transformers import Qwen2VLForConditionalGeneration, AutoProcessor, TrOCRProcessor, VisionEncoderDecoderModel
from qwen_vl_utils import process_vision_info
from paddleocr import PaddleOCR
import logging

# Setup logging
logging.basicConfig(level=logging.INFO, format='%(message)s')
logger = logging.getLogger(_name_)

# Document type categories for classification
DOCUMENT_TYPES = [
    "Engineering Drawing",
    "Maintenance Job Card",
    "Incident Report",
    "Vendor Invoice",
    "Purchase Order Correspondence",
    "Regulatory Directive",
    "Environmental Impact Study",
    "Safety Circular",
    "HR Policy",
    "Legal Opinion",
    "Board Meeting Minutes",
    "Other Document"
]


# ============================================================================
# QWEN-VL IMAGE DESCRIBER AND CLASSIFIER
# ============================================================================

class QwenImageDescriber:
    """Vision Language Model for image description and document classification"""
    
    def _init_(self, model_name="Qwen/Qwen2-VL-2B-Instruct", device=None):
        self.model_name = model_name
        
        # Auto-detect device
        if device is None:
            self.device = "cuda" if torch.cuda.is_available() else "cpu"
        else:
            self.device = device
            
        logger.info(f"Initializing Qwen-VL model: {model_name}")
        logger.info(f"Using device: {self.device}")
        
        try:
            self.model = Qwen2VLForConditionalGeneration.from_pretrained(
                model_name,
                torch_dtype=torch.float16 if self.device == "cuda" else torch.float32,
                device_map="auto" if self.device == "cuda" else None,
                resume_download=True,
                low_cpu_mem_usage=True
            )
            
            if self.device == "cpu":
                self.model = self.model.to(self.device)
            
            self.processor = AutoProcessor.from_pretrained(model_name, resume_download=True)
            
            logger.info("Qwen-VL model loaded successfully")
            self.available = True
            
        except Exception as e:
            logger.error(f"Failed to load Qwen-VL model: {str(e)}")
            self.available = False
            self.model = None
            self.processor = None
    
    def describe_image(self, image_path, prompt=None, max_new_tokens=512):
        """Generate detailed description of an image"""
        if not self.available:
            return "[ERROR] Qwen-VL model not available"
        
        if not os.path.exists(image_path):
            return f"[ERROR] Image file not found: {image_path}"
        
        try:
            if prompt is None:
                prompt = """Describe this image in detail. Include:
1. Main objects and subjects
2. Text visible in the image (if any)
3. Colors, layout, and composition
4. Context and scene description
5. Any notable details or features"""
            
            messages = [
                {
                    "role": "user",
                    "content": [
                        {"type": "image", "image": image_path},
                        {"type": "text", "text": prompt},
                    ],
                }
            ]
            
            text = self.processor.apply_chat_template(
                messages, tokenize=False, add_generation_prompt=True
            )
            
            image_inputs, video_inputs = process_vision_info(messages)
            
            inputs = self.processor(
                text=[text],
                images=image_inputs,
                videos=video_inputs,
                padding=True,
                return_tensors="pt",
            )
            
            inputs = inputs.to(self.device)
            
            with torch.no_grad():
                generated_ids = self.model.generate(
                    **inputs,
                    max_new_tokens=max_new_tokens,
                    do_sample=False
                )
            
            generated_ids_trimmed = [
                out_ids[len(in_ids):] for in_ids, out_ids in zip(inputs.input_ids, generated_ids)
            ]
            
            output_text = self.processor.batch_decode(
                generated_ids_trimmed,
                skip_special_tokens=True,
                clean_up_tokenization_spaces=False
            )[0]
            
            return output_text.strip()
            
        except Exception as e:
            logger.error(f"Error generating description for {image_path}: {str(e)}")
            return f"[ERROR] Failed to generate description: {str(e)}"
    
    def classify_document_type(self, image_path, max_new_tokens=256):
        """Classify document into predefined business categories"""
        prompt = """Analyze this document and classify it into ONE of these specific categories:

1. Engineering Drawing - Technical drawings, blueprints, CAD designs, schematics
2. Maintenance Job Card - Work orders, maintenance logs, service reports, repair records
3. Incident Report - Accident reports, safety incidents, occurrence records
4. Vendor Invoice - Bills, invoices, payment requests from suppliers
5. Purchase Order Correspondence - PO documents, procurement communications, order confirmations
6. Regulatory Directive - Government regulations, compliance directives, statutory orders
7. Environmental Impact Study - Environmental assessments, impact reports, sustainability studies
8. Safety Circular - Safety announcements, health & safety notices, warning bulletins
9. HR Policy - Human resources policies, employee handbooks, organizational policies
10. Legal Opinion - Legal advice, counsel opinions, legal memoranda
11. Board Meeting Minutes - Meeting minutes, board resolutions, corporate meeting records
12. Other Document - If none of the above categories fit

Provide your answer in this format:

DOCUMENT TYPE: [Selected category]
CONFIDENCE: [High/Medium/Low]
REASONING: [Brief explanation of why this classification was chosen based on visual elements, layout, content, headers, logos, or terminology visible in the document]

Be specific and look for key indicators like:
- Technical symbols/diagrams for engineering drawings
- Work order numbers/service dates for job cards
- Incident numbers/dates/signatures for incident reports
- Invoice numbers/amounts/vendor info for invoices
- PO numbers/item lists for purchase orders
- Official seals/legal language for regulatory directives
- Environmental data/impact assessments for environmental studies
- Safety symbols/warning language for safety circulars
- Policy headers/HR terminology for HR policies
- Legal formatting/opinions for legal documents
- Meeting dates/attendees/resolutions for meeting minutes"""
        
        return self.describe_image(image_path, prompt=prompt, max_new_tokens=max_new_tokens)


# ============================================================================
# OCR ENGINES INITIALIZATION
# ============================================================================

print("=" * 80)
print("Initializing OCR engines...")
print("=" * 80)

# PaddleOCR for English
ocr_en = None
try:
    ocr_en = PaddleOCR(lang='en', use_textline_orientation=False)
    print("  [OK] English PaddleOCR ready")
except Exception as e:
    print(f"  [WARNING] English OCR failed: {e}")

# TrOCR for handwriting
trocr_processor = None
trocr_model = None
try:
    trocr_processor = TrOCRProcessor.from_pretrained('microsoft/trocr-large-handwritten')
    trocr_model = VisionEncoderDecoderModel.from_pretrained('microsoft/trocr-large-handwritten')
    print("  [OK] TrOCR handwriting recognition ready")
except Exception as e:
    print(f"  [WARNING] TrOCR not available: {e}")

# Qwen-VL for image description and classification
qwen_describer = None
try:
    qwen_describer = QwenImageDescriber()
    if qwen_describer.available:
        print("  [OK] Qwen-VL image describer ready")
    else:
        qwen_describer = None
        print("  [WARNING] Qwen-VL model failed to initialize")
except Exception as e:
    print(f"  [WARNING] Qwen-VL initialization failed: {e}")
    qwen_describer = None

print("=" * 80)


# ============================================================================
# TEXT EXTRACTION FUNCTIONS
# ============================================================================

def extract_text_from_image(image_path):
    """Extract text from image using multi-engine OCR"""
    try:
        image = cv2.imread(str(image_path))
        if image is None:
            return ""
        
        all_texts = []
        
        # Try English OCR
        if ocr_en:
            try:
                result = ocr_en.ocr(image)
                if result and len(result) > 0 and result[0]:
                    ocr_result = result[0]
                    if hasattr(ocr_result, 'keys') and 'rec_texts' in ocr_result:
                        texts = [str(text).strip() for text in ocr_result['rec_texts'] if str(text).strip()]
                        if texts:
                            all_texts.extend(texts)
            except Exception as e:
                print(f"    PaddleOCR error: {str(e)[:80]}")
        
        # Try handwriting OCR if not much text found
        if trocr_processor and trocr_model and len(all_texts) < 3:
            try:
                pil_image = Image.open(image_path).convert("RGB")
                pixel_values = trocr_processor(pil_image, return_tensors="pt").pixel_values
                generated_ids = trocr_model.generate(pixel_values)
                text = trocr_processor.batch_decode(generated_ids, skip_special_tokens=True)[0]
                if text and text.strip():
                    all_texts.append(f"[Handwriting]: {text.strip()}")
            except Exception as e:
                print(f"    TrOCR error: {str(e)[:80]}")
        
        return "\n".join(all_texts) if all_texts else ""
    except Exception as e:
        print(f"    OCR Error: {str(e)[:100]}")
        return ""


def extract_text_from_pdf(pdf_path):
    """Extract text from PDF with OCR for scanned pages"""
    try:
        doc = fitz.open(pdf_path)
        all_text = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            
            if text.strip() and len(text.strip()) > 50:
                all_text.append(f"[Page {page_num + 1}]\n{text}")
            else:
                # Scanned page - use OCR
                print(f"    Page {page_num + 1}: Scanned page, applying OCR...")
                try:
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    img_bytes = pix.tobytes("png")
                    
                    # Save temporarily for OCR
                    temp_path = Path(f"temp_page_{page_num}.png")
                    with open(temp_path, 'wb') as f:
                        f.write(img_bytes)
                    
                    page_text = extract_text_from_image(temp_path)
                    
                    # Clean up
                    if temp_path.exists():
                        temp_path.unlink()
                    
                    if page_text:
                        all_text.append(f"[Page {page_num + 1} - OCR]\n{page_text}")
                except Exception as ocr_error:
                    print(f"    OCR failed for page {page_num + 1}: {str(ocr_error)[:80]}")
        
        doc.close()
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"    PDF error: {str(e)[:100]}")
        return ""


def extract_text_from_docx(docx_path):
    """Extract text from Word documents"""
    try:
        doc = docx_Document(docx_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs) if paragraphs else ""
    except Exception as e:
        print(f"    DOCX error: {str(e)[:100]}")
        return ""


def extract_text_from_pptx(pptx_path):
    """Extract text from PowerPoint presentations"""
    try:
        prs = Presentation(pptx_path)
        all_text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"=== Slide {slide_num} ==="]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:
                all_text.append("\n".join(slide_text))
        
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"    PPTX error: {str(e)[:100]}")
        return ""


def extract_text_from_xlsx(xlsx_path):
    """Extract text from Excel spreadsheets"""
    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
        all_text = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"=== Sheet: {sheet_name} ==="]
            
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    sheet_text.append(row_text)
            
            if len(sheet_text) > 1:
                all_text.append("\n".join(sheet_text))
        
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"    XLSX error: {str(e)[:100]}")
        return ""


def extract_text_from_txt(txt_path):
    """Read plain text file"""
    try:
        with open(txt_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        try:
            with open(txt_path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            print(f"    TXT error: {str(e)[:100]}")
            return ""


# ============================================================================
# DOCUMENT CLASSIFICATION FUNCTIONS
# ============================================================================

def classify_document(image_path):
    """Classify a document image into business categories"""
    if qwen_describer is None:
        return {
            "error": "Qwen-VL model not available",
            "document_type": "Unknown",
            "confidence": "N/A"
        }
    
    result_text = qwen_describer.classify_document_type(str(image_path))
    
    # Parse the result
    result = {
        "file": Path(image_path).name,
        "path": str(image_path),
        "document_type": "Unknown",
        "confidence": "Unknown",
        "reasoning": "",
        "raw_response": result_text
    }
    
    # Extract structured data from response
    lines = result_text.split('\n')
    for line in lines:
        line = line.strip()
        if line.startswith("DOCUMENT TYPE:"):
            result["document_type"] = line.replace("DOCUMENT TYPE:", "").strip()
        elif line.startswith("CONFIDENCE:"):
            result["confidence"] = line.replace("CONFIDENCE:", "").strip()
        elif line.startswith("REASONING:"):
            result["reasoning"] = line.replace("REASONING:", "").strip()
    
    return result


def generate_image_description(image_path):
    """Generate detailed description of image using Qwen-VL"""
    if qwen_describer is None:
        return None
    
    try:
        description = qwen_describer.describe_image(str(image_path))
        return description
    except Exception as e:
        print(f"    [WARNING] Qwen-VL description failed: {e}")
        return None


# ============================================================================
# MAIN PROCESSING FUNCTION
# ============================================================================

def process_file(input_path, output_dir):
    """Process a single file - extract text and classify if image"""
    file_path = Path(input_path)
    extension = file_path.suffix.lower()
    
    print(f"\n{'='*70}")
    print(f"Processing: {file_path.name}")
    print(f"Type: {extension}")
    print(f"{'='*70}")
    
    # Extract text based on file type
    text = ""
    
    try:
        if extension in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif', '.webp', '.gif']:
            text = extract_text_from_image(file_path)
        elif extension == '.pdf':
            text = extract_text_from_pdf(file_path)
        elif extension in ['.docx', '.doc']:
            text = extract_text_from_docx(file_path)
        elif extension in ['.pptx', '.ppt']:
            text = extract_text_from_pptx(file_path)
        elif extension in ['.xlsx', '.xls']:
            text = extract_text_from_xlsx(file_path)
        elif extension == '.txt':
            text = extract_text_from_txt(file_path)
        else:
            print(f"  [WARNING] Unsupported file type: {extension}")
            return False
    except KeyboardInterrupt:
        print(f"  [INTERRUPTED] Skipping...")
        return False
    except Exception as e:
        print(f"  [ERROR] Extraction failed: {str(e)[:100]}")
        text = f"[Error: {str(e)[:200]}]"
    
    # Save extracted text
    output_path = output_dir / f"{file_path.stem}.txt"
    
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            if text:
                f.write(text)
            else:
                f.write("[No text extracted]")
        
        char_count = len(text) if text else 0
        print(f"\n[SUCCESS] Saved to: {output_path.name}")
        print(f"          Characters: {char_count}")
        
        # For images: Generate description and classify document type
        if extension in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif', '.webp', '.gif'] and qwen_describer is not None:
            
            # Generate image description
            desc_path = output_dir / f"{file_path.stem}_description.txt"
            try:
                description = generate_image_description(file_path)
                if description and not description.startswith("[ERROR]"):
                    with open(desc_path, 'w', encoding='utf-8') as df:
                        df.write(f"Image Description (Qwen-VL)\n")
                        df.write(f"{'='*70}\n\n")
                        df.write(description)
                    print(f"          Image description saved to: {desc_path.name}")
            except Exception as e:
                print(f"          [WARNING] Description save failed: {str(e)[:80]}")
            
            # Classify document type
            classification_path = output_dir / f"{file_path.stem}_classification.txt"
            try:
                print(f"          Classifying document type...")
                classification_result = classify_document(file_path)
                
                if "error" not in classification_result:
                    with open(classification_path, 'w', encoding='utf-8') as cf:
                        cf.write(f"Document Type Classification (Qwen-VL)\n")
                        cf.write(f"{'='*70}\n\n")
                        cf.write(f"DOCUMENT TYPE: {classification_result['document_type']}\n")
                        cf.write(f"CONFIDENCE: {classification_result['confidence']}\n")
                        cf.write(f"REASONING: {classification_result['reasoning']}\n")
                    
                    print(f"          Document Type: {classification_result['document_type']}")
                    print(f"          Classification saved to: {classification_path.name}")
            except Exception as e:
                print(f"          [WARNING] Classification failed: {str(e)[:80]}")
        
        return True
    except Exception as e:
        print(f"  [ERROR] Failed to save: {str(e)[:100]}")
        return False


# ============================================================================
# MAIN FUNCTION
# ============================================================================

def main():
    """Process all files in input_files directory"""
    input_dir = Path("input_files")
    output_dir = Path("output")
    
    # Create output directory
    output_dir.mkdir(exist_ok=True)
    
    # Supported file extensions
    supported_extensions = [
        '.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif', '.webp', '.gif',  # Images
        '.pdf',  # PDF
        '.docx', '.doc',  # Word
        '.pptx', '.ppt',  # PowerPoint
        '.xlsx', '.xls',  # Excel
        '.txt'  # Text
    ]
    
    # Get all supported files
    files = [f for f in input_dir.iterdir() 
             if f.is_file() and f.suffix.lower() in supported_extensions]
    
    print(f"\n{'#'*70}")
    print(f"# Document Extraction and Classification System")
    print(f"# Total files found: {len(files)}")
    print(f"{'#'*70}\n")
    
    # Process each file
    success_count = 0
    fail_count = 0
    
    for idx, file in enumerate(files, 1):
        try:
            if process_file(file, output_dir):
                success_count += 1
            else:
                fail_count += 1
        except KeyboardInterrupt:
            print("\n[INTERRUPTED] Processing stopped by user")
            break
        except Exception as e:
            print(f"[ERROR] Unexpected error: {str(e)[:100]}")
            fail_count += 1

    # Summary
    print(f"\n{'='*70}")
    print(f"PROCESSING COMPLETE")
    print(f"{'='*70}")
    print(f"Total files: {len(files)}")
    print(f"Successful: {success_count}")
    print(f"Failed: {fail_count}")
    print(f"Output directory: {output_dir.absolute()}")
    print(f"{'='*70}\n")


if __name__ == "__main__":
    main()


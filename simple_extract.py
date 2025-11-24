"""
Simple text extraction script for all document types
Processes all files in input_files and saves extracted text to output folder
"""
import os
from pathlib import Path
from paddleocr import PaddleOCR
import cv2
import numpy as np
import fitz  # PyMuPDF
from docx import Document as docx_Document
from pptx import Presentation
import openpyxl

# Initialize PaddleOCR for both English and Malayalam
print("Initializing PaddleOCR for English and Malayalam...")
try:
    ocr_en = PaddleOCR(lang='en', use_textline_orientation=False)
    print("  English OCR ready")
except Exception as e:
    print(f"  Warning: English OCR failed: {e}")
    ocr_en = None

try:
    ocr_ml = PaddleOCR(lang='ml', use_textline_orientation=False)
    print("  Malayalam OCR ready")
except Exception as e:
    print(f"  Warning: Malayalam OCR not available: {e}")
    ocr_ml = None

# Try to initialize TrOCR for handwriting
try:
    from transformers import TrOCRProcessor, VisionEncoderDecoderModel
    from PIL import Image as PILImage
    trocr_processor = TrOCRProcessor.from_pretrained('microsoft/trocr-large-handwritten')
    trocr_model = VisionEncoderDecoderModel.from_pretrained('microsoft/trocr-large-handwritten')
    print("  Handwriting OCR (TrOCR) ready")
except Exception as e:
    print(f"  Warning: Handwriting OCR not available: {e}")
    trocr_processor = None
    trocr_model = None

def extract_text_from_image(image_path):
    """Extract text from image using PaddleOCR (both English and Malayalam)"""
    try:
        image = cv2.imread(str(image_path))
        if image is None:
            return ""
        
        all_texts = []
        
        # Try English OCR
        if ocr_en:
            try:
                result = ocr_en.ocr(image)
                if result and len(result) > 0 and result[0]:
                    ocr_result = result[0]
                    if hasattr(ocr_result, 'keys') and 'rec_texts' in ocr_result:
                        texts = [str(text).strip() for text in ocr_result['rec_texts'] if str(text).strip()]
                        if texts:
                            all_texts.extend(texts)
            except Exception as e:
                print(f"    English OCR Error: {str(e)[:80]}")
        
        # Try Malayalam OCR if English didn't find much text
        if ocr_ml and len(all_texts) < 3:
            try:
                result = ocr_ml.ocr(image)
                if result and len(result) > 0 and result[0]:
                    ocr_result = result[0]
                    if hasattr(ocr_result, 'keys') and 'rec_texts' in ocr_result:
                        texts = [str(text).strip() for text in ocr_result['rec_texts'] if str(text).strip()]
                        if texts:
                            all_texts.extend(texts)
            except Exception as e:
                print(f"    Malayalam OCR Error: {str(e)[:80]}")
        
        return "\n".join(all_texts) if all_texts else ""
    except KeyboardInterrupt:
        raise
    except Exception as e:
        print(f"    OCR Error: {str(e)[:100]}")
        return ""

def extract_text_from_pdf(pdf_path):
    """Extract text from PDF (text-based and scanned with OCR)"""
    try:
        doc = fitz.open(pdf_path)
        all_text = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Extract text directly
            text = page.get_text()
            
            if text.strip():
                all_text.append(text)
            else:
                # If no text, it's a scanned PDF - convert to image and use OCR
                print(f"    Page {page_num + 1}: Scanned page, using OCR...")
                try:
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom
                    img_bytes = pix.tobytes("png")
                    
                    # Convert to numpy array for OCR
                    import numpy as np
                    nparr = np.frombuffer(img_bytes, np.uint8)
                    img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                    
                    if img is not None:
                        # Try English OCR
                        page_text = []
                        if ocr_en:
                            result = ocr_en.ocr(img)
                            if result and len(result) > 0 and result[0]:
                                ocr_result = result[0]
                                if hasattr(ocr_result, 'keys') and 'rec_texts' in ocr_result:
                                    texts = [str(text).strip() for text in ocr_result['rec_texts']]
                                    page_text.extend(texts)
                        
                        # Try Malayalam if little text found
                        if ocr_ml and len(page_text) < 3:
                            result = ocr_ml.ocr(img)
                            if result and len(result) > 0 and result[0]:
                                ocr_result = result[0]
                                if hasattr(ocr_result, 'keys') and 'rec_texts' in ocr_result:
                                    texts = [str(text).strip() for text in ocr_result['rec_texts']]
                                    page_text.extend(texts)
                        
                        all_text.append("\n".join(page_text))
                except Exception as ocr_error:
                    print(f"    OCR failed for page {page_num + 1}: {str(ocr_error)[:80]}")
                    all_text.append(f"[Page {page_num + 1} - OCR failed]")
        
        doc.close()
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"    PDF Error: {str(e)[:100]}")
        return ""

def extract_text_from_docx(docx_path):
    """Extract text from Word documents with fallback options"""
    try:
        # Try python-docx first
        doc = docx_Document(docx_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        if paragraphs:
            return "\n".join(paragraphs)
    except Exception as e:
        print(f"    python-docx failed: {str(e)[:80]}")
    
    # Try docx2txt as fallback
    try:
        import docx2txt
        text = docx2txt.process(docx_path)
        if text and text.strip():
            return text
    except ImportError:
        print(f"    docx2txt not installed, skipping fallback")
    except Exception as e:
        print(f"    docx2txt failed: {str(e)[:80]}")
    
    # Last resort: Try to convert DOCX to image and OCR
    try:
        print(f"    Attempting OCR-based extraction...")
        import zipfile
        from PIL import Image
        import io
        
        with zipfile.ZipFile(docx_path) as doc_zip:
            # Look for images in the document
            image_files = [f for f in doc_zip.namelist() if f.startswith('word/media/')]
            
            all_text = []
            for img_file in image_files[:5]:  # Limit to first 5 images
                img_bytes = doc_zip.read(img_file)
                img = Image.open(io.BytesIO(img_bytes))
                img_cv = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
                
                # Try OCR on embedded images
                if ocr_en:
                    result = ocr_en.ocr(img_cv)
                    if result and len(result) > 0 and result[0]:
                        ocr_result = result[0]
                        if hasattr(ocr_result, 'keys') and 'rec_texts' in ocr_result:
                            texts = [str(text).strip() for text in ocr_result['rec_texts']]
                            all_text.extend(texts)
            
            if all_text:
                return "\n".join(all_text)
    except Exception as e:
        print(f"    OCR extraction failed: {str(e)[:80]}")
    
    return ""

def extract_text_from_pptx(pptx_path):
    """Extract text from PowerPoint"""
    try:
        prs = Presentation(pptx_path)
        all_text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            if slide_text:
                all_text.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))
        
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"Error processing PPTX {pptx_path}: {e}")
        return ""

def extract_text_from_xlsx(xlsx_path):
    """Extract text from Excel"""
    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
        all_text = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"Sheet: {sheet_name}"]
            
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    sheet_text.append(row_text)
            
            all_text.append("\n".join(sheet_text))
        
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"Error processing XLSX {xlsx_path}: {e}")
        return ""

def extract_text_from_txt(txt_path):
    """Read text file"""
    try:
        with open(txt_path, 'r', encoding='utf-8') as f:
            return f.read()
    except:
        try:
            with open(txt_path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            print(f"Error reading TXT {txt_path}: {e}")
            return ""

def process_file(input_path, output_dir):
    """Process a single file and extract text"""
    file_path = Path(input_path)
    extension = file_path.suffix.lower()
    
    print(f"\nProcessing: {file_path.name}")
    
    # Extract text based on file type
    text = ""
    
    try:
        if extension in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif']:
            text = extract_text_from_image(file_path)
        elif extension == '.pdf':
            text = extract_text_from_pdf(file_path)
        elif extension in ['.docx', '.doc']:
            text = extract_text_from_docx(file_path)
        elif extension in ['.pptx', '.ppt']:
            text = extract_text_from_pptx(file_path)
        elif extension in ['.xlsx', '.xls']:
            text = extract_text_from_xlsx(file_path)
        elif extension == '.txt':
            text = extract_text_from_txt(file_path)
        else:
            print(f"  Unsupported file type: {extension}")
            return False
    except KeyboardInterrupt:
        print(f"  [SKIP] Interrupted - skipping")
        return False
    except Exception as e:
        print(f"  [ERROR] Error: {str(e)[:100]}")
        text = f"[Error extracting text: {str(e)[:200]}]"
    
    # Save extracted text
    output_path = output_dir / f"{file_path.stem}.txt"
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text if text else "[No text extracted]")
        
        print(f"  [OK] Saved to: {output_path.name} ({len(text)} characters)")
        return True
    except Exception as e:
        print(f"  [ERROR] Failed to save: {e}")
        return False

def main():
    """Process all files in input_files directory"""
    input_dir = Path("input_files")
    output_dir = Path("output")
    
    # Create output directory if it doesn't exist
    output_dir.mkdir(exist_ok=True)
    
    # Supported file extensions
    supported_extensions = ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif', 
                          '.pdf', '.docx', '.doc', '.pptx', '.ppt', 
                          '.xlsx', '.xls', '.txt']
    
    # Get all supported files
    files = [f for f in input_dir.iterdir() 
             if f.is_file() and f.suffix.lower() in supported_extensions]
    
    print(f"\n{'='*60}")
    print(f"Found {len(files)} files to process")
    print(f"{'='*60}")
    
    successful = 0
    failed = 0
    
    for file_path in files:
        try:
            if process_file(file_path, output_dir):
                successful += 1
            else:
                failed += 1
        except Exception as e:
            print(f"  [ERROR] Error: {e}")
            failed += 1
    
    print(f"\n{'='*60}")
    print(f"Processing complete!")
    print(f"  [OK] Successful: {successful}")
    print(f"  [ERROR] Failed: {failed}")
    print(f"{'='*60}")

if __name__ == "__main__":
    main()

"""
Office Document Processor
Handles DOCX, XLSX, PPTX files using textract and specialized libraries
"""
import logging
from pathlib import Path
from typing import Dict, Any, List
import json

logger = logging.getLogger(__name__)


class OfficeDocumentProcessor:
    """Extract content from Office documents"""
    
    def __init__(self, config):
        self.config = config
        self._check_dependencies()
    
    def _check_dependencies(self):
        """Check availability of required libraries"""
        self.textract_available = False
        self.docx_available = False
        self.xlsx_available = False
        self.pptx_available = False
        
        try:
            import textract
            self.textract_available = True
            logger.info("textract available")
        except ImportError:
            logger.warning("textract not available")
        
        try:
            import docx
            self.docx_available = True
        except ImportError:
            logger.warning("python-docx not available")
        
        try:
            import openpyxl
            self.xlsx_available = True
        except ImportError:
            logger.warning("openpyxl not available")
        
        try:
            import pptx
            self.pptx_available = True
        except ImportError:
            logger.warning("python-pptx not available")
    
    def process_document(self, file_path: Path) -> Dict[str, Any]:
        """
        Process Office document and extract content
        
        Returns:
            Dict containing extracted content and metadata
        """
        suffix = file_path.suffix.lower()
        
        try:
            if suffix in ['.docx', '.doc']:
                return self.extract_docx(file_path)
            elif suffix in ['.xlsx', '.xls']:
                return self.extract_xlsx(file_path)
            elif suffix in ['.pptx', '.ppt']:
                return self.extract_pptx(file_path)
            else:
                return {
                    "text": "",
                    "success": False,
                    "error": f"Unsupported file type: {suffix}"
                }
        except Exception as e:
            logger.error(f"Office document processing failed: {e}")
            return {
                "text": "",
                "success": False,
                "error": str(e)
            }
    
    def extract_docx(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from DOCX file"""
        try:
            # Try textract first
            if self.textract_available:
                import textract
                text = textract.process(str(file_path)).decode('utf-8')
                return {
                    "text": text,
                    "method": "textract",
                    "success": True
                }
            
            # Fallback to python-docx
            if self.docx_available:
                import docx
                
                doc = docx.Document(file_path)
                
                content = {
                    "paragraphs": [],
                    "tables": [],
                    "headers_footers": []
                }
                
                # Extract paragraphs
                for para in doc.paragraphs:
                    if para.text.strip():
                        content["paragraphs"].append({
                            "text": para.text,
                            "style": para.style.name if para.style else "Normal"
                        })
                
                # Extract tables
                for table_idx, table in enumerate(doc.tables):
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    
                    content["tables"].append({
                        "index": table_idx,
                        "data": table_data
                    })
                
                # Combine all text
                full_text = "\n".join([p["text"] for p in content["paragraphs"]])
                
                # Add tables as text
                for table_info in content["tables"]:
                    full_text += f"\n\n[Table {table_info['index'] + 1}]\n"
                    for row in table_info["data"]:
                        full_text += " | ".join(row) + "\n"
                
                return {
                    "text": full_text,
                    "structured_content": content,
                    "method": "python-docx",
                    "success": True
                }
            
            return {
                "text": "",
                "success": False,
                "error": "No DOCX extraction library available"
            }
            
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            return {
                "text": "",
                "success": False,
                "error": str(e)
            }
    
    def extract_xlsx(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from XLSX file"""
        try:
            # Try textract first
            if self.textract_available:
                import textract
                text = textract.process(str(file_path)).decode('utf-8')
                return {
                    "text": text,
                    "method": "textract",
                    "success": True
                }
            
            # Fallback to openpyxl
            if self.xlsx_available:
                import openpyxl
                
                workbook = openpyxl.load_workbook(file_path, data_only=True)
                
                content = {
                    "sheets": []
                }
                
                full_text = ""
                
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    
                    sheet_data = []
                    for row in sheet.iter_rows(values_only=True):
                        # Filter out empty rows
                        if any(cell is not None for cell in row):
                            sheet_data.append([str(cell) if cell is not None else "" for cell in row])
                    
                    content["sheets"].append({
                        "name": sheet_name,
                        "data": sheet_data
                    })
                    
                    # Add to text
                    full_text += f"\n[Sheet: {sheet_name}]\n"
                    for row in sheet_data:
                        full_text += " | ".join(row) + "\n"
                
                return {
                    "text": full_text,
                    "structured_content": content,
                    "method": "openpyxl",
                    "success": True
                }
            
            return {
                "text": "",
                "success": False,
                "error": "No XLSX extraction library available"
            }
            
        except Exception as e:
            logger.error(f"XLSX extraction failed: {e}")
            return {
                "text": "",
                "success": False,
                "error": str(e)
            }
    
    def extract_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from PPTX file"""
        try:
            # Try textract first
            if self.textract_available:
                import textract
                text = textract.process(str(file_path)).decode('utf-8')
                return {
                    "text": text,
                    "method": "textract",
                    "success": True
                }
            
            # Fallback to python-pptx
            if self.pptx_available:
                from pptx import Presentation
                
                prs = Presentation(file_path)
                
                content = {
                    "slides": []
                }
                
                full_text = ""
                
                for slide_idx, slide in enumerate(prs.slides):
                    slide_content = {
                        "index": slide_idx,
                        "text_elements": [],
                        "tables": []
                    }
                    
                    slide_text = f"\n[Slide {slide_idx + 1}]\n"
                    
                    for shape in slide.shapes:
                        # Extract text from shapes
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_content["text_elements"].append(shape.text.strip())
                            slide_text += shape.text.strip() + "\n"
                        
                        # Extract tables
                        if shape.has_table:
                            table_data = []
                            for row in shape.table.rows:
                                row_data = [cell.text.strip() for cell in row.cells]
                                table_data.append(row_data)
                            
                            slide_content["tables"].append(table_data)
                            
                            # Add table to text
                            for row in table_data:
                                slide_text += " | ".join(row) + "\n"
                    
                    content["slides"].append(slide_content)
                    full_text += slide_text
                
                return {
                    "text": full_text,
                    "structured_content": content,
                    "method": "python-pptx",
                    "success": True
                }
            
            return {
                "text": "",
                "success": False,
                "error": "No PPTX extraction library available"
            }
            
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            return {
                "text": "",
                "success": False,
                "error": str(e)
            }

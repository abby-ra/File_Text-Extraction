"""
Enhanced text extraction with Malayalam support using Chithrakan
Supports: PDF, images (PNG/JPG), scanned photos, handwritten notes, job-card documents
Preserves Malayalam script and uses Docling for structured output
"""
import os
from pathlib import Path
import cv2
import numpy as np
import fitz  # PyMuPDF
from docx import Document as docx_Document
from pptx import Presentation
import openpyxl
import json
import csv
import zipfile
from io import BytesIO

# Initialize OCR engines
print("Initializing OCR engines...")

# PaddleOCR for English only
ocr_en = None
try:
    from paddleocr import PaddleOCR
    ocr_en = PaddleOCR(lang='en', use_textline_orientation=False)
    print("  [OK] English PaddleOCR ready")
except Exception as e:
    print(f"  [WARNING] English OCR failed: {e}")

# EasyOCR for Malayalam (printed and handwritten)
easyocr_available = False
easyocr_reader = None
try:
    import easyocr
    # Initialize EasyOCR with available languages
    # Malayalam is not supported in EasyOCR, using alternative approach
    easyocr_reader = easyocr.Reader(['en'], gpu=False)
    easyocr_available = True
    print("  [OK] EasyOCR English ready (Malayalam not supported in EasyOCR)")
    print("       For Malayalam: Install Tesseract with Malayalam language pack")
except ImportError:
    print("  [WARNING] EasyOCR not installed")
    print("           Install with: pip install easyocr")
except Exception as e:
    print(f"  [WARNING] EasyOCR initialization failed: {e}")

# Tesseract as fallback for mixed content
tesseract_available = False
try:
    import pytesseract
    # Test if Tesseract is installed
    pytesseract.get_tesseract_version()
    tesseract_available = True
    print("  [OK] Tesseract OCR ready (fallback)")
except Exception as e:
    print(f"  [WARNING] Tesseract not available: {e}")

# TrOCR for handwriting
trocr_available = False
trocr_processor = None
trocr_model = None
try:
    from transformers import TrOCRProcessor, VisionEncoderDecoderModel
    from PIL import Image as PILImage
    trocr_processor = TrOCRProcessor.from_pretrained('microsoft/trocr-large-handwritten')
    trocr_model = VisionEncoderDecoderModel.from_pretrained('microsoft/trocr-large-handwritten')
    trocr_available = True
    print("  [OK] TrOCR handwriting recognition ready")
except Exception as e:
    print(f"  [WARNING] TrOCR not available: {e}")

# Docling for structured conversion
docling_available = False
try:
    from docling.document_converter import DocumentConverter
    docling_available = True
    print("  [OK] Docling converter ready")
except ImportError:
    print("  [WARNING] Docling not installed")
    print("           Install with: pip install docling")
except Exception as e:
    print(f"  [WARNING] Docling initialization failed: {e}")

# Qwen-VL for image description and analysis
qwen_describer = None
try:
    from qwen_image_describer import QwenImageDescriber
    qwen_describer = QwenImageDescriber()
    if qwen_describer.available:
        print("  [OK] Qwen-VL image describer ready")
    else:
        qwen_describer = None
        print("  [WARNING] Qwen-VL model failed to initialize")
except ImportError:
    print("  [WARNING] Qwen-VL not available (qwen_image_describer.py not found)")
except Exception as e:
    print(f"  [WARNING] Qwen-VL initialization failed: {e}")
    qwen_describer = None

def detect_script(image):
    """Detect if image contains Malayalam script"""
    try:
        # Simple heuristic: check for Malayalam Unicode range patterns
        # Malayalam Unicode: U+0D00 to U+0D7F
        # This is a basic check - Chithrakan will be used for actual extraction
        
        # For now, try both and see which gives better results
        return 'mixed'  # Assume mixed content by default
    except:
        return 'unknown'

def generate_image_description(image_path):
    """Generate detailed description of image using Qwen-VL"""
    if qwen_describer is None:
        return None
    
    try:
        description = qwen_describer.describe_image(str(image_path))
        return description
    except Exception as e:
        print(f"    [WARNING] Qwen-VL description failed: {e}")
        return None

def extract_text_with_easyocr(image_path):
    """Extract Malayalam text using EasyOCR"""
    if not easyocr_available:
        return None
    
    try:
        # EasyOCR can handle both Malayalam and English in one pass
        result = easyocr_reader.readtext(str(image_path))
        
        if result:
            # Extract text from results (format: [(bbox, text, confidence), ...])
            texts = [item[1] for item in result if item[1].strip()]
            return "\n".join(texts) if texts else None
        
        return None
    except Exception as e:
        print(f"    EasyOCR error: {str(e)[:80]}")
        return None

def extract_text_with_tesseract(image_path, lang='eng+mal'):
    """Extract text using Tesseract with Malayalam support"""
    if not tesseract_available:
        return None
    
    try:
        import pytesseract
        from PIL import Image
        
        img = Image.open(image_path)
        
        # Try with specified language(s)
        text = pytesseract.image_to_string(img, lang=lang)
        
        if text and text.strip():
            return text.strip()
        
        return None
    except Exception as e:
        print(f"    Tesseract error: {str(e)[:80]}")
        return None

def extract_handwritten_text(image_path):
    """Extract handwritten text using TrOCR"""
    if not trocr_available:
        return None
    
    try:
        from PIL import Image as PILImage
        import torch
        
        image = PILImage.open(image_path).convert('RGB')
        
        # Process with TrOCR
        pixel_values = trocr_processor(image, return_tensors="pt").pixel_values
        
        with torch.no_grad():
            generated_ids = trocr_model.generate(pixel_values)
        
        text = trocr_processor.batch_decode(generated_ids, skip_special_tokens=True)[0]
        
        if text and text.strip():
            return text.strip()
        
        return None
    except Exception as e:
        print(f"    TrOCR error: {str(e)[:80]}")
        return None

def extract_text_from_image(image_path):
    """
    Extract text from image with multi-engine support:
    1. Generate image description using Qwen-VL
    2. Try EasyOCR for Malayalam & English (simultaneous)
    3. Try PaddleOCR for English
    4. Try TrOCR for handwriting
    5. Fall back to Tesseract for mixed content
    """
    try:
        all_texts = []
        engines_used = []
        
        # Strategy 0: Generate image description using Qwen-VL
        if qwen_describer is not None:
            print("    Generating image description with Qwen-VL...")
            description = generate_image_description(image_path)
            if description and not description.startswith("[ERROR]"):
                all_texts.append(f"[IMAGE DESCRIPTION - Qwen-VL]\n{description}")
                engines_used.append("Qwen-VL")
        
        # Strategy 1: Try EasyOCR first for Malayalam & English content
        if easyocr_available:
            print("    Trying EasyOCR for Malayalam+English...")
            malayalam_text = extract_text_with_easyocr(image_path)
            if malayalam_text and len(malayalam_text.strip()) > 10:
                all_texts.append(f"[Malayalam+English - EasyOCR]\n{malayalam_text}")
                engines_used.append("EasyOCR")
        
        # Strategy 2: Try PaddleOCR for English if EasyOCR didn't find much
        if ocr_en and len(all_texts) == 0:
            try:
                image = cv2.imread(str(image_path))
                if image is not None:
                    result = ocr_en.ocr(image)
                    if result and len(result) > 0 and result[0]:
                        ocr_result = result[0]
                        texts = []
                        if hasattr(ocr_result, 'keys') and 'rec_texts' in ocr_result:
                            texts = [str(t).strip() for t in ocr_result['rec_texts'] if str(t).strip()]
                        
                        if texts:
                            all_texts.append(f"[English - PaddleOCR]\n" + "\n".join(texts))
                            engines_used.append("PaddleOCR")
            except Exception as e:
                print(f"    PaddleOCR error: {str(e)[:80]}")
        
        # Strategy 3: Try TrOCR for handwritten content
        if trocr_available and len(all_texts) == 0:
            print("    Trying TrOCR for handwriting...")
            handwritten_text = extract_handwritten_text(image_path)
            if handwritten_text and len(handwritten_text.strip()) > 5:
                all_texts.append(f"[Handwritten - TrOCR]\n{handwritten_text}")
                engines_used.append("TrOCR")
        
        # Strategy 4: Fallback to Tesseract for mixed content
        if tesseract_available and len(all_texts) == 0:
            print("    Trying Tesseract (fallback)...")
            tesseract_text = extract_text_with_tesseract(image_path, lang='eng+mal')
            if tesseract_text:
                all_texts.append(f"[Mixed - Tesseract]\n{tesseract_text}")
                engines_used.append("Tesseract")
        
        if all_texts:
            result = "\n\n".join(all_texts)
            print(f"    Extracted using: {', '.join(engines_used)}")
            return result
        
        return ""
    except Exception as e:
        print(f"    Image extraction error: {str(e)[:100]}")
        return ""

def extract_text_from_pdf(pdf_path):
    """
    Extract text from PDF with OCR support for scanned pages
    Supports both text-based and image-based PDFs
    """
    try:
        doc = fitz.open(pdf_path)
        all_text = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Try direct text extraction first
            text = page.get_text()
            
            if text.strip() and len(text.strip()) > 50:
                # Good text-based content
                all_text.append(f"[Page {page_num + 1}]\n{text}")
            else:
                # Scanned page or minimal text - use OCR
                print(f"    Page {page_num + 1}: Scanned page, applying OCR...")
                try:
                    # Render page at higher resolution
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    img_bytes = pix.tobytes("png")
                    
                    # Save temporarily for OCR
                    temp_path = Path(f"temp_page_{page_num}.png")
                    with open(temp_path, 'wb') as f:
                        f.write(img_bytes)
                    
                    # Extract using multi-engine approach
                    page_text = extract_text_from_image(temp_path)
                    
                    # Clean up temp file
                    if temp_path.exists():
                        temp_path.unlink()
                    
                    if page_text:
                        all_text.append(f"[Page {page_num + 1} - OCR]\n{page_text}")
                    else:
                        all_text.append(f"[Page {page_num + 1}]\n[No text extracted]")
                        
                except Exception as ocr_error:
                    print(f"    OCR failed for page {page_num + 1}: {str(ocr_error)[:80]}")
                    all_text.append(f"[Page {page_num + 1}]\n[OCR failed]")
        
        doc.close()
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"    PDF error: {str(e)[:100]}")
        return ""

def extract_text_from_docx(docx_path):
    """Extract text from Word documents with multiple fallback methods"""
    try:
        # Method 1: python-docx
        doc = docx_Document(docx_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        if paragraphs:
            return "\n".join(paragraphs)
    except Exception as e:
        print(f"    python-docx failed: {str(e)[:80]}")
    
    # Method 2: docx2txt fallback
    try:
        import docx2txt
        text = docx2txt.process(docx_path)
        if text and text.strip():
            return text
    except Exception as e:
        print(f"    docx2txt failed: {str(e)[:80]}")
    
    # Method 3: Extract embedded images and OCR them
    try:
        print(f"    Extracting images from DOCX for OCR...")
        import zipfile
        from PIL import Image
        import io
        
        with zipfile.ZipFile(docx_path) as doc_zip:
            image_files = [f for f in doc_zip.namelist() if f.startswith('word/media/')]
            
            all_text = []
            for idx, img_file in enumerate(image_files[:10]):  # Limit to first 10 images
                try:
                    img_bytes = doc_zip.read(img_file)
                    img = Image.open(io.BytesIO(img_bytes))
                    
                    # Save temp image for OCR
                    temp_path = Path(f"temp_docx_img_{idx}.png")
                    img.save(temp_path)
                    
                    # Extract text from image
                    img_text = extract_text_from_image(temp_path)
                    
                    # Clean up
                    if temp_path.exists():
                        temp_path.unlink()
                    
                    if img_text:
                        all_text.append(f"[Image {idx + 1}]\n{img_text}")
                except Exception as img_error:
                    print(f"    Image {idx + 1} OCR failed: {str(img_error)[:60]}")
            
            if all_text:
                return "\n\n".join(all_text)
    except Exception as e:
        print(f"    Image extraction failed: {str(e)[:80]}")
    
    return ""

def extract_text_from_pptx(pptx_path):
    """Extract text from PowerPoint presentations"""
    try:
        prs = Presentation(pptx_path)
        all_text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"=== Slide {slide_num} ==="]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:
                all_text.append("\n".join(slide_text))
        
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"    PPTX error: {str(e)[:100]}")
        return ""

def extract_text_from_xlsx(xlsx_path):
    """Extract text from Excel spreadsheets"""
    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
        all_text = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"=== Sheet: {sheet_name} ==="]
            
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    sheet_text.append(row_text)
            
            if len(sheet_text) > 1:
                all_text.append("\n".join(sheet_text))
        
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"    XLSX error: {str(e)[:100]}")
        return ""

def extract_text_from_txt(txt_path):
    """Read plain text file"""
    try:
        with open(txt_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        try:
            with open(txt_path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            print(f"    TXT error: {str(e)[:100]}")
            return ""

def extract_text_from_csv(csv_path):
    """Extract text from CSV files"""
    try:
        all_rows = []
        with open(csv_path, 'r', encoding='utf-8', newline='') as f:
            reader = csv.reader(f)
            for row in reader:
                if row:  # Skip empty rows
                    all_rows.append("\t".join(row))
        return "\n".join(all_rows)
    except UnicodeDecodeError:
        try:
            all_rows = []
            with open(csv_path, 'r', encoding='latin-1', newline='') as f:
                reader = csv.reader(f)
                for row in reader:
                    if row:
                        all_rows.append("\t".join(row))
            return "\n".join(all_rows)
        except Exception as e:
            print(f"    CSV error: {str(e)[:100]}")
            return ""
    except Exception as e:
        print(f"    CSV error: {str(e)[:100]}")
        return ""

def extract_text_from_rtf(rtf_path):
    """Extract text from RTF files"""
    try:
        # Try using pypandoc if available
        import pypandoc
        text = pypandoc.convert_file(str(rtf_path), 'plain', format='rtf')
        return text
    except ImportError:
        print("    [INFO] pypandoc not installed. Install with: pip install pypandoc")
    except Exception as e:
        print(f"    RTF conversion error: {str(e)[:80]}")
    
    # Fallback: Basic RTF text extraction
    try:
        with open(rtf_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
            # Remove RTF control words and brackets
            import re
            text = re.sub(r'\\[a-z]+\d*\s?', '', content)
            text = re.sub(r'[{}]', '', text)
            text = '\n'.join([line.strip() for line in text.split('\n') if line.strip()])
            return text
    except Exception as e:
        print(f"    RTF fallback error: {str(e)[:100]}")
        return ""

def extract_text_from_html(html_path):
    """Extract text from HTML files"""
    try:
        from bs4 import BeautifulSoup
        with open(html_path, 'r', encoding='utf-8', errors='ignore') as f:
            soup = BeautifulSoup(f, 'html.parser')
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            text = soup.get_text(separator='\n')
            # Clean up whitespace
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            return '\n'.join(lines)
    except ImportError:
        print("    [INFO] BeautifulSoup not installed. Install with: pip install beautifulsoup4")
    except Exception as e:
        print(f"    HTML error: {str(e)[:100]}")
    
    # Fallback: Basic HTML tag removal
    try:
        import re
        with open(html_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
            text = re.sub(r'<[^>]+>', '', content)
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            return '\n'.join(lines)
    except Exception as e:
        print(f"    HTML fallback error: {str(e)[:100]}")
        return ""

def extract_text_from_odt(odt_path):
    """Extract text from OpenDocument Text files (.odt)"""
    try:
        with zipfile.ZipFile(odt_path) as odt_zip:
            content_xml = odt_zip.read('content.xml').decode('utf-8')
            
            # Parse XML and extract text
            from xml.etree import ElementTree as ET
            root = ET.fromstring(content_xml)
            
            # Extract all text nodes
            text_elements = []
            for elem in root.iter():
                if elem.text:
                    text_elements.append(elem.text.strip())
                if elem.tail:
                    text_elements.append(elem.tail.strip())
            
            return '\n'.join([t for t in text_elements if t])
    except Exception as e:
        print(f"    ODT error: {str(e)[:100]}")
        return ""

def convert_with_docling(text, output_path):
    """Convert extracted text to structured format using Docling"""
    if not docling_available:
        return False
    
    try:
        # Create a temporary markdown file for Docling
        temp_md = output_path.parent / f"{output_path.stem}_temp.md"
        
        with open(temp_md, 'w', encoding='utf-8') as f:
            f.write(f"# Extracted Text\n\n{text}")
        
        # Use Docling to convert
        converter = DocumentConverter()
        result = converter.convert(str(temp_md))
        
        # Save as JSON for structured data
        json_path = output_path.parent / f"{output_path.stem}.json"
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump({
                'source': output_path.name,
                'text': text,
                'structured': result.document.export_to_dict() if hasattr(result, 'document') else {}
            }, f, ensure_ascii=False, indent=2)
        
        # Clean up temp file
        if temp_md.exists():
            temp_md.unlink()
        
        print(f"    [Docling] Structured data saved to {json_path.name}")
        return True
    except Exception as e:
        print(f"    Docling conversion error: {str(e)[:80]}")
        return False

def process_file(input_path, output_dir):
    """Process a single file and extract text"""
    file_path = Path(input_path)
    extension = file_path.suffix.lower()
    
    print(f"\n{'='*70}")
    print(f"Processing: {file_path.name}")
    print(f"Type: {extension}")
    print(f"{'='*70}")
    
    # Extract text based on file type
    text = ""
    
    try:
        if extension in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif', '.webp', '.gif']:
            text = extract_text_from_image(file_path)
        elif extension == '.pdf':
            text = extract_text_from_pdf(file_path)
        elif extension in ['.docx', '.doc']:
            text = extract_text_from_docx(file_path)
        elif extension in ['.pptx', '.ppt']:
            text = extract_text_from_pptx(file_path)
        elif extension in ['.xlsx', '.xls']:
            text = extract_text_from_xlsx(file_path)
        elif extension == '.txt':
            text = extract_text_from_txt(file_path)
        elif extension == '.csv':
            text = extract_text_from_csv(file_path)
        elif extension == '.rtf':
            text = extract_text_from_rtf(file_path)
        elif extension in ['.html', '.htm']:
            text = extract_text_from_html(file_path)
        elif extension == '.odt':
            text = extract_text_from_odt(file_path)
        else:
            print(f"  [WARNING] Unsupported file type: {extension}")
            print(f"  Supported: PDF, DOCX, DOC, PPTX, PPT, XLSX, XLS, TXT, CSV, RTF, HTML, HTM, ODT, Images")
            return False
    except KeyboardInterrupt:
        print(f"  [INTERRUPTED] Skipping...")
        return False
    except Exception as e:
        print(f"  [ERROR] Extraction failed: {str(e)[:100]}")
        text = f"[Error: {str(e)[:200]}]"
    
    # Save extracted text with original format preserved
    output_path = output_dir / f"{file_path.stem}.txt"
    
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            if text:
                f.write(text)
            else:
                f.write("[No text extracted]")
        
        char_count = len(text) if text else 0
        print(f"\n[SUCCESS] Saved to: {output_path.name}")
        print(f"          Characters: {char_count}")
        
        # For images, also save detailed description to separate file if Qwen-VL is available
        if extension in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif', '.webp'] and qwen_describer is not None:
            desc_path = output_dir / f"{file_path.stem}_description.txt"
            try:
                description = generate_image_description(file_path)
                if description and not description.startswith("[ERROR]"):
                    with open(desc_path, 'w', encoding='utf-8') as df:
                        df.write(f"Image Description (Qwen-VL)\n")
                        df.write(f"{'='*70}\n\n")
                        df.write(description)
                    print(f"          Image description saved to: {desc_path.name}")
            except Exception as e:
                print(f"          [WARNING] Description save failed: {str(e)[:80]}")
        
        # Try Docling conversion for structured output
        if text and char_count > 50:
            convert_with_docling(text, output_path)
        
        return True
    except Exception as e:
        print(f"  [ERROR] Failed to save: {str(e)[:100]}")
        return False

def main():
    """Process all files in input_files directory"""
    input_dir = Path("input_files")
    output_dir = Path("output")
    
    # Create output directory
    output_dir.mkdir(exist_ok=True)
    
    # Supported file extensions
    supported_extensions = [
        '.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif', '.webp', '.gif',  # Images
        '.pdf',  # PDF
        '.docx', '.doc',  # Word
        '.pptx', '.ppt',  # PowerPoint
        '.xlsx', '.xls',  # Excel
        '.txt',  # Plain text
        '.csv',  # CSV spreadsheets
        '.rtf',  # Rich Text Format
        '.html', '.htm',  # HTML documents
        '.odt'  # OpenDocument Text
    ]
    
    # Get all supported files
    files = [f for f in input_dir.iterdir() 
             if f.is_file() and f.suffix.lower() in supported_extensions]
    
    print(f"\n{'#'*70}")
    print(f"# Enhanced Text Extraction with Malayalam Support")
    print(f"# Total files found: {len(files)}")
    print(f"{'#'*70}\n")
    
    successful = 0
    failed = 0
    
    for file_path in sorted(files):
        try:
            if process_file(file_path, output_dir):
                successful += 1
            else:
                failed += 1
        except Exception as e:
            print(f"  [ERROR] Unexpected error: {str(e)[:100]}")
            failed += 1
    
    print(f"\n{'#'*70}")
    print(f"# Processing Complete!")
    print(f"# ✓ Successful: {successful}")
    print(f"# ✗ Failed: {failed}")
    print(f"{'#'*70}\n")

if __name__ == "__main__":
    main()

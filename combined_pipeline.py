"""
Combined single-file entrypoint for text extraction.

Usage:
  python combined_pipeline.py --input_dir input_files --output_dir extracted_texts --summarize

This inlines the Mineru client and the IntelligentTextExtractor from the repo.
"""
import os
import sys
import logging
import json
from pathlib import Path
from typing import List, Dict, Optional

# Ensure Tesseract (Windows) path extension if present in original project
os.environ['PATH'] += os.pathsep + r'C:\Program Files\Tesseract-OCR'

logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

try:
    from PIL import Image
except Exception:
    Image = None

try:
    import pytesseract
except Exception:
    pytesseract = None

# Optional image & ML utilities (imported lazily where used)
try:
    import cv2
except Exception:
    cv2 = None

try:
    import numpy as np
except Exception:
    np = None

try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    import camelot
except Exception:
    camelot = None

try:
    # sklearn fallback classifier
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.linear_model import LogisticRegression
    SKLEARN_AVAILABLE = True
except Exception:
    TfidfVectorizer = None
    LogisticRegression = None
    SKLEARN_AVAILABLE = False

try:
    # transformers-based classifier (optional)
    from transformers import AutoTokenizer, AutoModelForSequenceClassification
    TRANSFORMERS_AVAILABLE = True
except Exception:
    AutoTokenizer = None
    AutoModelForSequenceClassification = None
    TRANSFORMERS_AVAILABLE = False

try:
    import magic
except Exception:
    magic = None

# Defer boto3 import to runtime to avoid heavy import cost during module import.
boto3 = None

try:
    import requests
except Exception:
    requests = None

import io


class MineruClient:
    """Minimal client wrapper for a Mineru-like extraction REST API."""
    def __init__(self, api_url: Optional[str] = None, api_key: Optional[str] = None, timeout: int = 60):
        self.api_url = api_url or os.environ.get('MINERU_API_URL', 'https://api.mineru.ai/v1/extract')
        self.api_key = api_key or os.environ.get('MINERU_API_KEY')
        self.timeout = timeout

    def is_configured(self) -> bool:
        return bool(self.api_key)

    def extract(self, file_path: str) -> Dict:
        if not requests:
            logger.warning("requests not installed; MineruClient unavailable")
            return {}

        if not os.path.exists(file_path):
            logger.error("Mineru extract failed: file not found %s", file_path)
            return {}

        headers = {}
        if self.api_key:
            headers['Authorization'] = f"Bearer {self.api_key}"

        try:
            with open(file_path, 'rb') as f:
                files = {'file': (os.path.basename(file_path), f)}
                resp = requests.post(self.api_url, headers=headers, files=files, timeout=self.timeout)

            if resp.status_code != 200:
                logger.error("Mineru API error (%s): %s", resp.status_code, resp.text)
                return {}

            try:
                data = resp.json()
            except Exception as e:
                logger.error("Failed to decode JSON from Mineru response: %s", e)
                return {}

            if isinstance(data, dict):
                return data
            else:
                logger.error("Unexpected Mineru response format: %s", type(data))
                return {}

        except Exception as e:
            logger.exception("Mineru extraction failed: %s", e)
            return {}


class IntelligentTextExtractor:
    def __init__(self, input_folder: str = "input_files", output_folder: str = "extracted_texts"):
        self.input_folder = input_folder
        self.output_folder = output_folder
        self.setup_folders()

        # Textract client (optional). Import boto3 only if explicitly enabled to avoid
        # importing heavy AWS libraries during normal module import.
        if os.environ.get('ENABLE_AWS_TEXTRACT', '0') == '1':
            try:
                import boto3 as _boto3
                self.textract_client = _boto3.client('textract', region_name='us-east-1')
                logger.info('AWS Textract client initialized')
            except Exception as e:
                logger.warning(f'AWS Textract not configured: {e}')
                self.textract_client = None
        else:
            self.textract_client = None

        # Mineru client
        try:
            self.mineru_client = MineruClient()
            if not self.mineru_client.is_configured():
                logger.info("Mineru client not configured (no MINERU_API_KEY). Skipping Mineru integration.")
                self.mineru_client = None
            else:
                logger.info("Mineru client configured. Remote extraction enabled.")
        except Exception as e:
            logger.warning(f"Failed to initialize Mineru client: {e}")
            self.mineru_client = None

    def setup_folders(self):
        os.makedirs(self.input_folder, exist_ok=True)
        os.makedirs(self.output_folder, exist_ok=True)
        logger.info(f"Created folders: {self.input_folder}, {self.output_folder}")

    def detect_file_type(self, file_path: str) -> str:
        try:
            if magic:
                mime_type = magic.from_file(file_path, mime=True)
                logger.info(f"Detected MIME type for {file_path}: {mime_type}")
                if mime_type == 'application/pdf':
                    return 'pdf'
                elif mime_type.startswith('image/'):
                    return 'image'
                elif mime_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
                    return 'docx'
                elif mime_type.startswith('text/'):
                    return 'text'
                else:
                    return 'unknown'
        except Exception:
            pass

        # Fallback to extension
        ext = Path(file_path).suffix.lower()
        if ext == '.pdf':
            return 'pdf'
        elif ext in ['.jpg', '.jpeg', '.png', '.tiff', '.bmp']:
            return 'image'
        elif ext in ['.docx', '.doc']:
            return 'docx'
        elif ext in ['.txt', '.text']:
            return 'text'
        return 'unknown'

    # -------------------------
    # Image preprocessing helpers
    # -------------------------
    def preprocess_image(self, image_path: str, deskew: bool = True, binarize: bool = True, denoise: bool = True, resize_max: int = 2000) -> Optional[Path]:
        """Return a Path to a preprocessed temporary image suitable for OCR.
        If OpenCV/Pillow are not available, returns the original path."""
        src = Path(image_path)
        if not src.exists():
            return None

        if cv2 is None or np is None or Image is None:
            return src

        try:
            # Read with OpenCV
            img = cv2.imdecode(np.fromfile(str(src), dtype=np.uint8), cv2.IMREAD_COLOR)
            if img is None:
                return src

            # Resize if very large
            h, w = img.shape[:2]
            max_dim = max(h, w)
            if resize_max and max_dim > resize_max:
                scale = resize_max / float(max_dim)
                img = cv2.resize(img, (int(w * scale), int(h * scale)), interpolation=cv2.INTER_AREA)

            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

            if denoise:
                try:
                    gray = cv2.fastNlMeansDenoising(gray, None, 10, 7, 21)
                except Exception:
                    pass

            if deskew:
                try:
                    coords = np.column_stack(np.where(gray > 0))
                    angle = cv2.minAreaRect(coords)[-1]
                    if angle < -45:
                        angle = -(90 + angle)
                    else:
                        angle = -angle
                    (h, w) = gray.shape[:2]
                    center = (w // 2, h // 2)
                    M = cv2.getRotationMatrix2D(center, angle, 1.0)
                    gray = cv2.warpAffine(gray, M, (w, h), flags=cv2.INTER_CUBIC, borderMode=cv2.BORDER_REPLICATE)
                except Exception:
                    pass

            if binarize:
                try:
                    gray = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2)
                except Exception:
                    _, gray = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)

            out = Path("temp_preprocessed.png")
            # Write using numpy-aware method to support unicode paths
            ext = cv2.imencode('.png', gray)[1]
            ext.tofile(str(out))
            return out
        except Exception:
            return src

    def analyze_pdf_content(self, pdf_path: str) -> Dict:
        if not fitz:
            return {'has_text': False, 'has_images': False, 'has_tables': False, 'total_pages': 0, 'text_pages': [], 'image_pages': [], 'table_pages': []}

        doc = fitz.open(pdf_path)
        analysis = {'has_text': False, 'has_images': False, 'has_tables': False, 'total_pages': len(doc), 'text_pages': [], 'image_pages': [], 'table_pages': []}
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text().strip()
            if text:
                analysis['has_text'] = True
                analysis['text_pages'].append(page_num)
            if page.get_images():
                analysis['has_images'] = True
                analysis['image_pages'].append(page_num)
            if text and ('\t' in text or text.count('|') > 5):
                analysis['has_tables'] = True
                analysis['table_pages'].append(page_num)
        doc.close()
        logger.info(f"PDF analysis: {analysis}")
        return analysis

    def extract_text_from_pdf(self, pdf_path: str) -> str:
        if not fitz:
            logger.warning("PyMuPDF not installed; cannot extract PDF text")
            return ""
        doc = fitz.open(pdf_path)
        text_content = []
        logger.info(f"Extracting text from {len(doc)} pages using PyMuPDF...")
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            if text.strip():
                text_content.append(f"--- Page {page_num + 1} ---\n{text}")
            if (page_num + 1) % 50 == 0:
                logger.info(f"Processed {page_num + 1}/{len(doc)} pages")
        doc.close()
        result = "\n\n".join(text_content)
        if not result.strip():
            logger.info("No selectable text found, using OCR fallback")
            # Attempt table-aware extraction first (pdfplumber/camelot) if available
            table_text = ""
            try:
                if pdfplumber is not None:
                    try:
                        with pdfplumber.open(pdf_path) as pp:
                            tables_found = 0
                            for p in pp.pages:
                                for tbl in p.extract_tables():
                                    tables_found += 1
                                    # format table as TSV-like text
                                    for row in tbl:
                                        table_text += "\t".join([str(c) if c is not None else "" for c in row]) + "\n"
                            if tables_found:
                                logger.info(f"Extracted {tables_found} tables via pdfplumber")
                    except Exception:
                        table_text = ""

                # Camelot provides better structure for some PDFs (requires ghostscript)
                if not table_text and camelot is not None:
                    try:
                        cams = camelot.read_pdf(str(pdf_path), pages='all')
                        if cams and len(cams) > 0:
                            for c in cams:
                                table_text += c.df.to_csv(sep='\t', index=False)
                            logger.info(f"Extracted {len(cams)} tables via camelot")
                    except Exception:
                        table_text = ""
            except Exception:
                table_text = ""

            if table_text:
                result = table_text
            else:
                result = self.extract_images_from_pdf_ocr(pdf_path)
        return result

    def extract_images_from_pdf_ocr(self, pdf_path: str, page_numbers: List[int] = None) -> str:
        if not fitz or not Image or not pytesseract:
            logger.warning("OCR prerequisites missing")
            return ""

        doc = fitz.open(pdf_path)
        ocr_content = []
        pages_to_process = page_numbers or range(len(doc))
        logger.info(f"Extracting text from {len(pages_to_process)} pages using OCR...")

        for page_num in pages_to_process:
            page = doc[page_num]
            # try higher resolution first, fallback on MemoryError to lower resolution
            tried_scales = [2.0, 1.5, 1.0]
            img_data = None
            for scale in tried_scales:
                try:
                    mat = fitz.Matrix(scale, scale)
                    pix = page.get_pixmap(matrix=mat)
                    img_data = pix.tobytes("png")
                    break
                except MemoryError:
                    logger.warning(f"MemoryError rendering page {page_num + 1} at scale {scale}, trying lower scale")
                    continue
                except Exception as e:
                    logger.error(f"Failed to render page {page_num + 1}: {e}")
                    img_data = None
                    break

            if not img_data:
                logger.error(f"Skipping page {page_num + 1} due to render failure")
                continue

            try:
                image = Image.open(io.BytesIO(img_data))
                ocr_text = pytesseract.image_to_string(image, lang='mal+eng') if pytesseract else ''
                if ocr_text.strip():
                    ocr_content.append(f"--- Page {page_num + 1} (OCR) ---\n{ocr_text}")
            except Exception as e:
                logger.error(f"OCR failed for page {page_num + 1}: {e}")

            if (page_num + 1) % 10 == 0:
                logger.info(f"OCR processed {page_num + 1} pages")

        doc.close()
        return "\n\n".join(ocr_content)

    def extract_tables_from_pdf_textract(self, pdf_path: str, page_numbers: List[int] = None) -> str:
        if not self.textract_client:
            logger.warning("Textract not available, skipping table extraction")
            return ""
        try:
            with open(pdf_path, 'rb') as document:
                response = self.textract_client.analyze_document(Document={'Bytes': document.read()}, FeatureTypes=['TABLES'])
            table_content = []
            for block in response.get('Blocks', []):
                if block.get('BlockType') == 'TABLE':
                    table_text = self.parse_textract_table(block, response.get('Blocks', []))
                    table_content.append(table_text)
            return "\n\n".join(table_content)
        except Exception as e:
            logger.error(f"Textract extraction failed: {e}")
            return ""

    def parse_textract_table(self, table_block: Dict, all_blocks: List[Dict]) -> str:
        table_text = "--- TABLE ---\n"
        for relationship in table_block.get('Relationships', []):
            if relationship.get('Type') == 'CHILD':
                for child_id in relationship.get('Ids', []):
                    child_block = next((b for b in all_blocks if b.get('Id') == child_id), None)
                    if child_block and child_block.get('BlockType') == 'CELL':
                        if 'Text' in child_block:
                            table_text += child_block['Text'] + " | "
        return table_text

    def extract_from_image(self, image_path: str) -> str:
        if not Image or not pytesseract:
            logger.warning("PIL/pytesseract not installed; skipping image OCR")
            return ""
        try:
            logger.info(f"Extracting text from image: {image_path}")
            image = Image.open(image_path)
            text = pytesseract.image_to_string(image, lang='mal+eng')
            return text
        except Exception as e:
            logger.error(f"Image OCR failed: {e}")
            return ""

    def extract_from_docx(self, docx_path: str) -> str:
        try:
            from docx import Document
        except Exception:
            logger.warning("python-docx not installed; skipping DOCX extraction")
            return ""
        try:
            logger.info(f"Extracting text from Word document: {docx_path}")
            doc = Document(docx_path)
            text_content = [p.text for p in doc.paragraphs if p.text.strip()]
            result = "\n".join(text_content)
            return result
        except Exception as e:
            logger.error(f"Word document extraction failed: {e}")
            return ""

    def extract_from_text_file(self, text_path: str) -> str:
        try:
            logger.info(f"Reading text file: {text_path}")
            with open(text_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            logger.error(f"Text file reading failed: {e}")
            return ""

    def process_single_file(self, file_path: str) -> str:
        file_type = self.detect_file_type(file_path)
        logger.info(f"Processing {file_path} as {file_type}")
        all_text = []

        if self.mineru_client:
            try:
                logger.info("Attempting extraction via Mineru API for %s", file_path)
                mineru_result = self.mineru_client.extract(file_path)
                if mineru_result and mineru_result.get('text'):
                    all_text.append("=== MINERU EXTRACTED ===\n" + mineru_result.get('text'))
                    if 'tables' in mineru_result and mineru_result['tables']:
                        all_text.append("=== MINERU TABLES ===\n" + json.dumps(mineru_result['tables'], ensure_ascii=False))
                    logger.info("Mineru extraction succeeded for %s", file_path)
                    return "\n\n".join(all_text)
                else:
                    logger.info("Mineru returned no text for %s, falling back to local extraction", file_path)
            except Exception as e:
                logger.warning(f"Mineru extraction failed for {file_path}: {e}. Falling back to local extraction.")

        if file_type == 'pdf':
            analysis = self.analyze_pdf_content(file_path)
            if analysis['has_text']:
                text_content = self.extract_text_from_pdf(file_path)
                if text_content.strip():
                    all_text.append("=== TEXT CONTENT ===\n" + text_content)
            if analysis['has_images']:
                ocr_content = self.extract_images_from_pdf_ocr(file_path, analysis['image_pages'])
                if ocr_content.strip():
                    all_text.append("=== OCR CONTENT ===\n" + ocr_content)
            if analysis['has_tables'] and self.textract_client:
                table_content = self.extract_tables_from_pdf_textract(file_path, analysis['table_pages'])
                if table_content.strip():
                    all_text.append("=== TABLE CONTENT ===\n" + table_content)

        elif file_type == 'image':
            image_text = self.extract_from_image(file_path)
            if image_text.strip():
                all_text.append(image_text)

        elif file_type == 'docx':
            docx_text = self.extract_from_docx(file_path)
            if docx_text.strip():
                all_text.append(docx_text)

        elif file_type == 'text':
            text_content = self.extract_from_text_file(file_path)
            if text_content.strip():
                all_text.append(text_content)

        else:
            logger.warning(f"Unsupported file type: {file_type}")
            return ""

        return "\n\n".join(all_text)

    def save_extracted_text(self, text: str, original_filename: str) -> str:
        output_filename = f"{Path(original_filename).stem}_extracted.txt"
        output_path = os.path.join(self.output_folder, output_filename)
        try:
            with open(output_path, 'w', encoding='utf-8') as file:
                file.write(text)
            logger.info(f"Saved extracted text to: {output_path}")
            return output_path
        except Exception as e:
            logger.error(f"Failed to save text: {e}")
            return ""

    def process_all_files(self):
        if not os.path.exists(self.input_folder):
            logger.error(f"Input folder {self.input_folder} does not exist")
            return
        files = [f for f in os.listdir(self.input_folder) if os.path.isfile(os.path.join(self.input_folder, f))]
        if not files:
            logger.info(f"No files found in {self.input_folder}")
            return
        logger.info(f"Found {len(files)} files to process")
        for filename in files:
            file_path = os.path.join(self.input_folder, filename)
            logger.info(f"Processing file: {filename}")
            try:
                extracted_text = self.process_single_file(file_path)
                if extracted_text.strip():
                    output_path = self.save_extracted_text(extracted_text, filename)
                    logger.info(f"Successfully processed {filename} -> {output_path}")
                else:
                    logger.warning(f"No text extracted from {filename}")
            except Exception as e:
                logger.error(f"Failed to process {filename}: {e}")


def run_full_pipeline(input_dir: str = "input_files", output_dir: str = "extracted_texts", call_summarizer: bool = True) -> int:
    """Run extraction for all files in `input_dir` and save to `output_dir`.
    If `call_summarizer` is True and `hugging_face.py` exists, it will be invoked as a subprocess.
    Returns 0 on success, non-zero on failure."""
    extractor = IntelligentTextExtractor(input_folder=input_dir, output_folder=output_dir)

    # Check tesseract availability (best-effort)
    if pytesseract:
        try:
            pytesseract.get_tesseract_version()
            logger.info("Tesseract OCR is available")
        except Exception as e:
            logger.error(f"Tesseract OCR not found: {e}")
    else:
        logger.info("pytesseract not installed; OCR will be skipped where required")

    extractor.process_all_files()

    if call_summarizer:
        hugging_path = os.path.join(os.getcwd(), 'hugging_face.py')
        if os.path.exists(hugging_path):
            logger.info("Invoking summarizer: hugging_face.py")
            import subprocess
            try:
                subprocess.run([sys.executable, hugging_path], check=True)
                logger.info("Summarizer completed")
            except Exception as e:
                logger.warning(f"Summarizer failed: {e}")
        else:
            logger.info("No hugging_face.py found; skipping summarization")

    return 0


def _cli():
    import argparse
    parser = argparse.ArgumentParser(description='Combined text extraction pipeline')
    parser.add_argument('--input_dir', default='input_files')
    parser.add_argument('--output_dir', default='extracted_texts')
    parser.add_argument('--deepdoctection', action='store_true', help='Run DeepDoctection engine (image description & classification)')
    parser.add_argument('--deep-input', default='input_files', help='Input directory for deepdoctection (defaults to input_files)')
    parser.add_argument('--deep-output', default='output', help='Output directory for deepdoctection (defaults to output)')
    parser.add_argument('--no-extra', action='store_true', help='Do not write auxiliary outputs (classification, descriptions). Only write extracted (and summarized if enabled)')
    parser.add_argument('--no-summarize', action='store_true', help='Do not call hugging_face.py')
    args = parser.parse_args()
    # Control whether to write auxiliary/metadata outputs
    global SKIP_EXTRA_OUTPUTS
    SKIP_EXTRA_OUTPUTS = bool(args.no_extra)
    if args.deepdoctection:
        # Run the deep document analysis engine which handles images and classification
        try:
            deepdoctection_main(input_dir=args.deep_input, output_dir=args.deep_output)
            rc = 0
        except Exception as e:
            logger.error(f"DeepDoctection run failed: {e}")
            rc = 2
    else:
        rc = run_full_pipeline(input_dir=args.input_dir, output_dir=args.output_dir, call_summarizer=not args.no_summarize)
    sys.exit(rc)


# ---------------------------------------------------------------------------
# Inlined DeepDoctectionAnalyzer (adapted to avoid heavy imports/initialization at import time)
# ---------------------------------------------------------------------------

try:
    import cv2
except Exception:
    cv2 = None

try:
    import numpy as np
except Exception:
    np = None

try:
    from docx import Document as docx_Document
except Exception:
    docx_Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    import torch
except Exception:
    torch = None

try:
    from transformers import Qwen2VLForConditionalGeneration, AutoProcessor, TrOCRProcessor, VisionEncoderDecoderModel
except Exception:
    Qwen2VLForConditionalGeneration = None
    AutoProcessor = None
    TrOCRProcessor = None
    VisionEncoderDecoderModel = None

try:
    from qwen_vl_utils import process_vision_info
except Exception:
    process_vision_info = None

try:
    from paddleocr import PaddleOCR
except Exception:
    PaddleOCR = None

try:
    from PIL import Image as PIL_Image
except Exception:
    PIL_Image = None

import logging as _logging
_logging.basicConfig(level=_logging.INFO, format='%(message)s')
_logger = _logging.getLogger(__name__)

DOCUMENT_TYPES = [
    "Engineering Drawing",
    "Maintenance Job Card",
    "Incident Report",
    "Vendor Invoice",
    "Purchase Order Correspondence",
    "Regulatory Directive",
    "Environmental Impact Study",
    "Safety Circular",
    "HR Policy",
    "Legal Opinion",
    "Board Meeting Minutes",
    "Other Document"
]


class TextClassifier:
    """Flexible classifier for document types.
    Tries transformers-based classifier if available; falls back to sklearn TF-IDF + LogisticRegression;
    finally falls back to keyword rules."""
    def __init__(self, labels=None):
        self.labels = labels or DOCUMENT_TYPES
        self.mode = 'keyword'
        self._sk_vectorizer = None
        self._sk_model = None
        self._tf_tokenizer = None
        self._tf_model = None

        if TRANSFORMERS_AVAILABLE:
            try:
                # default small model — user can change
                self._tf_tokenizer = AutoTokenizer.from_pretrained('distilbert-base-uncased')
                # user may fine-tune and point to model path in future
                self._tf_model = None
                self.mode = 'transformers'
            except Exception:
                self._tf_tokenizer = None
                self._tf_model = None

        if not TRANSFORMERS_AVAILABLE and SKLEARN_AVAILABLE:
            try:
                # create an empty TF-IDF + LR model placeholder; user should train and load in production
                self._sk_vectorizer = TfidfVectorizer(max_features=20000)
                self._sk_model = LogisticRegression(max_iter=1000)
                self.mode = 'sklearn'
            except Exception:
                self._sk_vectorizer = None
                self._sk_model = None

        # simple keyword map as last resort
        self.keyword_map = {
            'invoice': ['invoice', 'amount due', 'total due', 'invoice number', 'bill to'],
            'purchase order': ['purchase order', 'po number', 'purchase order no'],
            'maintenance': ['maintenance', 'job card', 'work order', 'technician'],
            'incident report': ['incident report', 'incident', 'reported', 'injury'],
            'board meeting': ['board meeting', 'minutes', 'attendees', 'resolution'],
            'legal opinion': ['legal opinion', 'opinion of counsel', 'whereas', 'hereto'],
            'hr policy': ['hr policy', 'human resources', 'leave policy', 'code of conduct']
        }

    def predict(self, text: str) -> Dict:
        if not text or not text.strip():
            return {'label': 'Unknown', 'confidence': 0.0}

        # Transformers path (if a fine-tuned model is present this would be used)
        if self.mode == 'transformers' and self._tf_tokenizer is not None and self._tf_model is not None:
            try:
                inputs = self._tf_tokenizer(text, truncation=True, padding=True, return_tensors='pt')
                outputs = self._tf_model(**inputs)
                probs = outputs.logits.softmax(-1).detach().cpu().numpy()[0]
                idx = int(probs.argmax())
                return {'label': self.labels[idx] if idx < len(self.labels) else 'Unknown', 'confidence': float(probs[idx])}
            except Exception:
                pass

        # Sklearn fallback (if model has been trained and set)
        if self.mode == 'sklearn' and self._sk_vectorizer is not None and self._sk_model is not None:
            try:
                x = self._sk_vectorizer.transform([text])
                probs = self._sk_model.predict_proba(x)[0]
                idx = int(probs.argmax())
                return {'label': self.labels[idx] if idx < len(self.labels) else 'Unknown', 'confidence': float(probs[idx])}
            except Exception:
                pass

        # Keyword heuristic last-resort
        text_l = text.lower()
        scores = {}
        for key, kws in self.keyword_map.items():
            for kw in kws:
                if kw in text_l:
                    scores[key] = scores.get(key, 0) + 1
        if scores:
            best = max(scores.items(), key=lambda x: x[1])
            return {'label': best[0], 'confidence': float(best[1]) / (len(text.split()) + 1)}

        return {'label': 'Unknown', 'confidence': 0.0}


class QwenImageDescriber:
    """Vision Language Model wrapper (lightly adapted)."""
    def __init__(self, model_name="Qwen/Qwen2-VL-2B-Instruct", device=None):
        self.model_name = model_name
        if device is None:
            self.device = "cuda" if (torch is not None and torch.cuda.is_available()) else "cpu"
        else:
            self.device = device

        _logger.info(f"Initializing Qwen-VL model: {model_name}")
        _logger.info(f"Using device: {self.device}")

        try:
            if Qwen2VLForConditionalGeneration is None or AutoProcessor is None:
                raise RuntimeError("Transformers Qwen-VL classes not available in environment")

            self.model = Qwen2VLForConditionalGeneration.from_pretrained(
                model_name,
                torch_dtype=torch.float16 if self.device == "cuda" else torch.float32,
                device_map="auto" if self.device == "cuda" else None,
                resume_download=True,
                low_cpu_mem_usage=True
            )

            if self.device == "cpu":
                self.model = self.model.to(self.device)

            self.processor = AutoProcessor.from_pretrained(model_name, resume_download=True)
            _logger.info("Qwen-VL model loaded successfully")
            self.available = True

        except Exception as e:
            _logger.error(f"Failed to load Qwen-VL model: {str(e)}")
            self.available = False
            self.model = None
            self.processor = None

    def describe_image(self, image_path, prompt=None, max_new_tokens=512):
        # Lightweight, defensive wrapper to generate an image description using Qwen-VL.
        if not self.available:
            return "[ERROR] Qwen-VL model not available"
        if not os.path.exists(image_path):
            return f"[ERROR] Image file not found: {image_path}"

        if prompt is None:
            prompt = (
                "Describe this image in detail. Include main objects, text visible, colors, layout, "
                "context and any notable details."
            )

        messages = [
            {"role": "user", "content": [{"type": "image", "image": image_path}, {"type": "text", "text": prompt}]}
        ]

        if process_vision_info is None:
            return "[ERROR] qwen_vl_utils.process_vision_info not available"

        try:
            # Attempt to build a text prompt using the processor if available.
            text = prompt
            if hasattr(self.processor, 'apply_chat_template'):
                try:
                    text = self.processor.apply_chat_template(messages, tokenize=False, add_generation_prompt=True)
                except Exception:
                    text = prompt

            image_inputs, video_inputs = process_vision_info(messages)

            inputs = self.processor(text=[text], images=image_inputs, videos=video_inputs, padding=True, return_tensors="pt")

            # Move inputs to device where possible (may not be supported for some processor outputs)
            try:
                inputs = inputs.to(self.device)
            except Exception:
                pass

            with torch.no_grad():
                generated_ids = self.model.generate(**inputs, max_new_tokens=max_new_tokens, do_sample=False)

            # Try trimming prompt tokens before decoding; fallback to raw decode on failure
            try:
                generated_ids_trimmed = [out_ids[len(in_ids):] for in_ids, out_ids in zip(inputs.input_ids, generated_ids)]
                output_text = self.processor.batch_decode(generated_ids_trimmed, skip_special_tokens=True, clean_up_tokenization_spaces=False)[0]
            except Exception:
                output_text = self.processor.batch_decode(generated_ids, skip_special_tokens=True, clean_up_tokenization_spaces=False)[0]

            return output_text.strip()
        except Exception as e:
            _logger.error(f"Error generating description for {image_path}: {str(e)}")
            return f"[ERROR] Failed to generate description: {str(e)}"

    def classify_document_type(self, image_path, max_new_tokens=256):
        prompt = """Analyze this document and classify it into ONE of the specified categories."""
        return self.describe_image(image_path, prompt=prompt, max_new_tokens=max_new_tokens)


# NOTE: The original DeepDoctectionAnalyzer executed heavy model initialization at import
# time. To keep the combined file safe to import, we provide an explicit initializer.

class DeepDocEngine:
    def __init__(self):
        self.ocr_en = None
        self.trocr_processor = None
        self.trocr_model = None
        self.qwen_describer = None

    def initialize_engines(self):
        # PaddleOCR
        if PaddleOCR is not None:
            try:
                self.ocr_en = PaddleOCR(lang='en', use_textline_orientation=False)
                print("  [OK] English PaddleOCR ready")
            except Exception as e:
                print(f"  [WARNING] English OCR failed: {e}")

        # TrOCR
        if TrOCRProcessor is not None and VisionEncoderDecoderModel is not None:
            try:
                self.trocr_processor = TrOCRProcessor.from_pretrained('microsoft/trocr-large-handwritten')
                self.trocr_model = VisionEncoderDecoderModel.from_pretrained('microsoft/trocr-large-handwritten')
                print("  [OK] TrOCR handwriting recognition ready")
            except Exception as e:
                print(f"  [WARNING] TrOCR not available: {e}")

        # Qwen-VL
        try:
            if Qwen2VLForConditionalGeneration is not None:
                self.qwen_describer = QwenImageDescriber()
                if self.qwen_describer.available:
                    print("  [OK] Qwen-VL image describer ready")
                else:
                    self.qwen_describer = None
                    print("  [WARNING] Qwen-VL model failed to initialize")
        except Exception as e:
            print(f"  [WARNING] Qwen-VL initialization failed: {e}")
            self.qwen_describer = None

    # The following methods are thin wrappers adapted from the original file
    def extract_text_from_image(self, image_path):
        try:
            if cv2 is None:
                return ""
            image = cv2.imread(str(image_path))
            if image is None:
                return ""
            all_texts = []
            if self.ocr_en:
                try:
                    result = self.ocr_en.ocr(image)
                    if result and len(result) > 0 and result[0]:
                        ocr_result = result[0]
                        if hasattr(ocr_result, 'keys') and 'rec_texts' in ocr_result:
                            texts = [str(text).strip() for text in ocr_result['rec_texts'] if str(text).strip()]
                            if texts:
                                all_texts.extend(texts)
                except Exception as e:
                    print(f"    PaddleOCR error: {str(e)[:80]}")

            if self.trocr_processor and self.trocr_model and len(all_texts) < 3:
                try:
                    pil_image = PIL_Image.open(image_path).convert("RGB")
                    pixel_values = self.trocr_processor(pil_image, return_tensors="pt").pixel_values
                    generated_ids = self.trocr_model.generate(pixel_values)
                    text = self.trocr_processor.batch_decode(generated_ids, skip_special_tokens=True)[0]
                    if text and text.strip():
                        all_texts.append(f"[Handwriting]: {text.strip()}")
                except Exception as e:
                    print(f"    TrOCR error: {str(e)[:80]}")

            return "\n".join(all_texts) if all_texts else ""
        except Exception as e:
            print(f"    OCR Error: {str(e)[:100]}")
            return ""

    def extract_text_from_pdf(self, pdf_path):
        if fitz is None:
            return ""
        try:
            doc = fitz.open(pdf_path)
            all_text = []
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text()
                if text.strip() and len(text.strip()) > 50:
                    all_text.append(f"[Page {page_num + 1}]\n{text}")
                else:
                    try:
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        img_bytes = pix.tobytes("png")
                        temp_path = Path(f"temp_page_{page_num}.png")
                        with open(temp_path, 'wb') as f:
                            f.write(img_bytes)
                        page_text = self.extract_text_from_image(temp_path)
                        if temp_path.exists():
                            temp_path.unlink()
                        if page_text:
                            all_text.append(f"[Page {page_num + 1} - OCR]\n{page_text}")
                    except Exception as ocr_error:
                        print(f"    OCR failed for page {page_num + 1}: {str(ocr_error)[:80]}")
            doc.close()
            return "\n\n".join(all_text)
        except Exception as e:
            print(f"    PDF error: {str(e)[:100]}")
            return ""

    def extract_text_from_docx(self, docx_path):
        if docx_Document is None:
            return ""
        try:
            doc = docx_Document(docx_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs) if paragraphs else ""
        except Exception as e:
            print(f"    DOCX error: {str(e)[:100]}")
            return ""

    def extract_text_from_pptx(self, pptx_path):
        if Presentation is None:
            return ""
        try:
            prs = Presentation(pptx_path)
            all_text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_lines = [f"=== Slide {slide_num} ==="]
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text and shape.text.strip():
                            slide_lines.append(shape.text)
                    except Exception:
                        # ignore shapes that cause issues
                        continue
                if len(slide_lines) > 1:
                    all_text.append("\n".join(slide_lines))
            return "\n\n".join(all_text)
        except Exception as e:
            print(f"    PPTX error: {str(e)[:100]}")
            return ""

    def extract_text_from_xlsx(self, xlsx_path):
        if openpyxl is None:
            return ""
        try:
            wb = openpyxl.load_workbook(xlsx_path, data_only=True)
            all_text = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"=== Sheet: {sheet_name} ==="]
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        sheet_text.append(row_text)
                if len(sheet_text) > 1:
                    all_text.append("\n".join(sheet_text))
            return "\n\n".join(all_text)
        except Exception as e:
            print(f"    XLSX error: {str(e)[:100]}")
            return ""

    def extract_text_from_txt(self, txt_path):
        try:
            with open(txt_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            try:
                with open(txt_path, 'r', encoding='latin-1') as f:
                    return f.read()
            except Exception as e:
                print(f"    TXT error: {str(e)[:100]}")
                return ""

    def classify_document(self, image_path):
        if self.qwen_describer is None:
            return {"error": "Qwen-VL model not available", "document_type": "Unknown", "confidence": "N/A"}
        result_text = self.qwen_describer.classify_document_type(str(image_path))
        result = {"file": Path(image_path).name, "path": str(image_path), "document_type": "Unknown", "confidence": "Unknown", "reasoning": "", "raw_response": result_text}
        lines = result_text.split('\n')
        for line in lines:
            line = line.strip()
            if line.startswith("DOCUMENT TYPE:"):
                result["document_type"] = line.replace("DOCUMENT TYPE:", "").strip()
            elif line.startswith("CONFIDENCE:"):
                result["confidence"] = line.replace("CONFIDENCE:", "").strip()
            elif line.startswith("REASONING:"):
                result["reasoning"] = line.replace("REASONING:", "").strip()
        return result

    def generate_image_description(self, image_path):
        if self.qwen_describer is None:
            return None
        try:
            description = self.qwen_describer.describe_image(str(image_path))
            return description
        except Exception as e:
            print(f"    [WARNING] Qwen-VL description failed: {e}")
            return None


    def process_file(self, input_path, output_dir: Path):
        file_path = Path(input_path)
        extension = file_path.suffix.lower()
        print(f"\n{'='*70}")
        print(f"Processing: {file_path.name}")
        print(f"Type: {extension}")
        print(f"{'='*70}")
        text = ""
        try:
            if extension in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif', '.webp', '.gif']:
                # Preprocess image (deskew / binarize) before OCR for better results
                try:
                    pre = self.preprocess_image(str(file_path))
                    if pre and pre.exists():
                        text = self.extract_text_from_image(pre)
                    else:
                        text = self.extract_text_from_image(file_path)
                except Exception:
                    text = self.extract_text_from_image(file_path)
            elif extension == '.pdf':
                text = self.extract_text_from_pdf(file_path)
            elif extension in ['.docx', '.doc']:
                text = self.extract_text_from_docx(file_path)
            elif extension in ['.pptx', '.ppt']:
                text = self.extract_text_from_pptx(file_path)
            elif extension in ['.xlsx', '.xls']:
                text = self.extract_text_from_xlsx(file_path)
            elif extension == '.txt':
                text = self.extract_text_from_txt(file_path)
            else:
                print(f"  [WARNING] Unsupported file type: {extension}")
                return False
        except KeyboardInterrupt:
            print(f"  [INTERRUPTED] Skipping...")
            return False
        except Exception as e:
            print(f"  [ERROR] Extraction failed: {str(e)[:100]}")
            text = f"[Error: {str(e)[:200]}]"

        output_path = output_dir / f"{file_path.stem}.txt"
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                if text:
                    f.write(text)
                else:
                    f.write("[No text extracted]")
            char_count = len(text) if text else 0
            print(f"\n[SUCCESS] Saved to: {output_path.name}")
            print(f"          Characters: {char_count}")
            # Auxiliary outputs (classification, descriptions) removed by user request.
            # Only the extracted text file is written by default.
            return True
        except Exception as e:
            print(f"  [ERROR] Failed to save: {str(e)[:100]}")
            return False


def deepdoctection_main(input_dir: Optional[str] = None, output_dir: Optional[str] = None):
    input_dir = Path(input_dir or "input_files")
    output_dir = Path(output_dir or "output")
    output_dir.mkdir(exist_ok=True)
    supported_extensions = ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif', '.webp', '.gif', '.pdf', '.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls', '.txt']
    files = [f for f in input_dir.iterdir() if f.is_file() and f.suffix.lower() in supported_extensions]
    print(f"\n{'#'*70}")
    print(f"# Document Extraction and Classification System")
    print(f"# Total files found: {len(files)}")
    print(f"{'#'*70}\n")
    engine = DeepDocEngine()
    engine.initialize_engines()
    success_count = 0
    fail_count = 0
    for file in files:
        try:
            if engine.process_file(file, output_dir):
                success_count += 1
            else:
                fail_count += 1
        except KeyboardInterrupt:
            print("\n[INTERRUPTED] Processing stopped by user")
            break
        except Exception as e:
            print(f"[ERROR] Unexpected error: {str(e)[:100]}")
            fail_count += 1
    print(f"\n{'='*70}")
    print(f"PROCESSING COMPLETE")
    print(f"{'='*70}")
    print(f"Total files: {len(files)}")
    print(f"Successful: {success_count}")
    print(f"Failed: {fail_count}")
    print(f"Output directory: {output_dir.absolute()}")


if __name__ == '__main__':
    _cli()

